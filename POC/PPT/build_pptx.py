"""
Elyra — Investor Pitch Deck (PPTX generator, v2)
Strict layout grid · refined glassmorphism · Final-POC parity.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# DESIGN TOKENS  (mirrored from ultimate.css :root)
# ---------------------------------------------------------------------------
NIGHT       = RGBColor(0x06, 0x08, 0x0D)
INK         = RGBColor(0x0D, 0x13, 0x20)
PANEL       = RGBColor(0x10, 0x17, 0x25)
PANEL_2     = RGBColor(0x16, 0x1F, 0x31)
PANEL_3     = RGBColor(0x1C, 0x27, 0x3D)
PANEL_TOP   = RGBColor(0x22, 0x2E, 0x46)   # for inner highlight
STROKE      = RGBColor(0x33, 0x40, 0x57)
STROKE_SOFT = RGBColor(0x22, 0x2C, 0x40)
STROKE_HI   = RGBColor(0x46, 0x55, 0x6F)

CREAM       = RGBColor(0xFB, 0xF7, 0xEF)
GOLD        = RGBColor(0xFF, 0xC5, 0x6E)
CORAL       = RGBColor(0xFF, 0x6F, 0x8E)
TEAL        = RGBColor(0x43, 0xE2, 0xD0)
VIOLET      = RGBColor(0x93, 0x84, 0xFF)
BLUE        = RGBColor(0x6E, 0xA8, 0xFF)
GREEN       = RGBColor(0x71, 0xE6, 0xA4)

SOFT        = RGBColor(0xC8, 0xD0, 0xDA)
MUTED       = RGBColor(0x88, 0x93, 0xA3)
FAINT       = RGBColor(0x59, 0x64, 0x77)
INK_TEXT    = RGBColor(0x16, 0x10, 0x08)

FONT_DISPLAY = "Georgia"
FONT_BODY    = "Segoe UI"
FONT_MONO    = "Consolas"

# ---------------------------------------------------------------------------
# GRID (single source of truth — every slide obeys this)
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

MARGIN_L   = Inches(0.75)
MARGIN_R   = Inches(0.75)
CONTENT_W  = SW - MARGIN_L - MARGIN_R                # 11.83"

KICKER_Y   = Inches(0.62)
KICKER_H   = Inches(0.28)
H2_Y       = Inches(0.96)
H2_H_2L    = Inches(1.55)                            # two-line headline
H2_H_1L    = Inches(0.92)                            # single-line
LEDE_Y     = Inches(2.55)
LEDE_H     = Inches(0.78)
DIV_Y      = Inches(3.40)                            # accent divider
BODY_Y     = Inches(3.62)
BODY_H     = Inches(3.30)
FOOTER_DIV = Inches(7.04)
FOOTER_Y   = Inches(7.10)


# ---------------------------------------------------------------------------
# LOW-LEVEL XML / FILL HELPERS
# ---------------------------------------------------------------------------
def _strip_effects(shp):
    sppr = shp._element.spPr
    for ext in sppr.findall(qn('a:effectLst')):
        sppr.remove(ext)
    etree.SubElement(sppr, qn('a:effectLst'))


def shape_alpha(shape, opacity_pct):
    sf = shape.fill._xPr.find(qn('a:solidFill'))
    if sf is None:
        return
    clr = sf.find(qn('a:srgbClr'))
    if clr is None:
        clr = sf.find(qn('a:schemeClr'))
    if clr is None:
        return
    for existing in clr.findall(qn('a:alpha')):
        clr.remove(existing)
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(opacity_pct * 1000)))


def line_alpha(shape, opacity_pct):
    ln = shape.line._get_or_add_ln()
    sf = ln.find(qn('a:solidFill'))
    if sf is None:
        return
    clr = sf.find(qn('a:srgbClr'))
    if clr is None:
        return
    for existing in clr.findall(qn('a:alpha')):
        clr.remove(existing)
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(opacity_pct * 1000)))


def add_rect(slide, x, y, w, h, fill=PANEL, line=None, line_w=0.75,
             radius=None, alpha=None, line_alpha_pct=None):
    if radius is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if alpha is not None:
        shape_alpha(shp, alpha)

    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if line_alpha_pct is not None:
            line_alpha(shp, line_alpha_pct)

    _strip_effects(shp)
    return shp


def add_ellipse(slide, x, y, w, h, fill, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if alpha is not None:
        shape_alpha(shp, alpha)
    shp.line.fill.background()
    _strip_effects(shp)
    return shp


def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14,
             color=CREAM, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, tracking=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    runs = [(text, {})] if isinstance(text, str) else text
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, (t, opts) in enumerate(runs):
        run = p.add_run()
        run.text = t
        f = run.font
        f.name = opts.get('font', font)
        f.size = Pt(opts.get('size', size))
        f.bold = opts.get('bold', bold)
        f.italic = opts.get('italic', italic)
        f.color.rgb = opts.get('color', color)
        if 'tracking' in opts or tracking is not None:
            tr = opts.get('tracking', tracking)
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(int(tr)))
    return tb


# ---------------------------------------------------------------------------
# AMBIENT LAYERS
# ---------------------------------------------------------------------------
def add_bg(slide, *, orbs=('coral', 'teal'), grid=True, vignette=True):
    # base night
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NIGHT
    bg.line.fill.background(); _strip_effects(bg)

    # ink tint plate (pseudo-gradient)
    plate = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    plate.fill.solid(); plate.fill.fore_color.rgb = INK
    shape_alpha(plate, 30)
    plate.line.fill.background(); _strip_effects(plate)

    # subtle dot-grid via fine lines (soft 56-px style cell)
    if grid:
        gw = RGBColor(0xFF, 0xFF, 0xFF)
        step = Inches(0.74)
        x = step
        while x < SW:
            ln = slide.shapes.add_connector(1, x, 0, x, SH)
            ln.line.color.rgb = gw; ln.line.width = Pt(0.4)
            line_alpha(ln, 4)
            x += step
        y = step
        while y < SH:
            ln = slide.shapes.add_connector(1, 0, y, SW, y)
            ln.line.color.rgb = gw; ln.line.width = Pt(0.4)
            line_alpha(ln, 4)
            y += step

    # orb specs (bigger, much softer, parked at edges)
    orb_specs = {
        'coral':  (CORAL,  Inches( 9.4), Inches(-3.4), Inches(7.8), 14),
        'teal':   (TEAL,   Inches(-3.6), Inches( 2.4), Inches(7.4), 13),
        'gold':   (GOLD,   Inches( 8.0), Inches( 4.6), Inches(6.6), 11),
        'violet': (VIOLET, Inches(-1.6), Inches( 4.8), Inches(5.6),  9),
        'blue':   (BLUE,   Inches(-2.4), Inches(-2.6), Inches(6.0), 12),
    }
    for o in orbs:
        if o in orb_specs:
            c, x, y, d, a = orb_specs[o]
            add_ellipse(slide, x, y, d, d, c, alpha=a)

    # corner vignette plates (adds depth around edges)
    if vignette:
        v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        v.fill.solid(); v.fill.fore_color.rgb = NIGHT
        shape_alpha(v, 12); v.line.fill.background(); _strip_effects(v)


# ---------------------------------------------------------------------------
# COMPONENTS
# ---------------------------------------------------------------------------
def glass_card(slide, x, y, w, h, *, radius=0.06, fill=PANEL,
               alpha=85, line=STROKE, line_alpha_pct=70, line_w=0.75,
               highlight=True):
    """A glassmorphic card: dark fill + thin stroke + 1-pt top inner highlight."""
    card = add_rect(slide, x, y, w, h,
                    fill=fill, line=line, line_w=line_w,
                    radius=radius, alpha=alpha,
                    line_alpha_pct=line_alpha_pct)
    if highlight:
        # 1-pt bright inner highlight at top, simulating glass top edge
        hl = slide.shapes.add_connector(1, x + Inches(0.06), y + Inches(0.02),
                                        x + w - Inches(0.06), y + Inches(0.02))
        hl.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hl.line.width = Pt(0.6)
        line_alpha(hl, 16)
    return card


def brand_mark(slide, x, y, scale=1.0, sub="Pehchaan Layer · India-first"):
    d = int(Inches(0.55) * scale)
    add_ellipse(slide, x, y, d, d, GOLD)
    add_ellipse(slide, x, y, d, d, CORAL, alpha=55)
    add_text(slide, x, y, d, d, 'E', font=FONT_DISPLAY, size=int(20 * scale),
             color=INK_TEXT, bold=True, italic=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = x + d + Inches(0.16)
    add_text(slide, tx, y - Inches(0.02), Inches(3.6), Inches(0.30),
             'Elyra', font=FONT_BODY, size=int(15 * scale),
             color=CREAM, bold=True)
    add_text(slide, tx, y + Inches(0.26), Inches(4.0), Inches(0.24),
             sub, font=FONT_BODY, size=int(8.5 * scale),
             color=MUTED, bold=True, tracking=80)


def kicker(slide, x, y, w, text, color=GOLD, size=10):
    add_text(slide, x, y, w, KICKER_H, text.upper(),
             font=FONT_BODY, size=size, color=color,
             bold=True, tracking=240)


def chip(slide, x, y, w, h, text, *, fill_color=PANEL_2, text_color=SOFT,
         dot_color=None, border=STROKE, size=9.5, mono=False,
         center=False, alpha=88):
    add_rect(slide, x, y, w, h, fill=fill_color, line=border,
             line_w=0.6, radius=0.5, alpha=alpha, line_alpha_pct=70)
    inner_x = x + Inches(0.18)
    if dot_color is not None:
        d = Inches(0.12); dy = y + (h - d) / 2
        add_ellipse(slide, inner_x, dy, d, d, dot_color)
        text_x = inner_x + d + Inches(0.10); text_w = w - (text_x - x) - Inches(0.18)
    else:
        text_x = inner_x; text_w = w - Inches(0.36)
    add_text(slide, text_x, y, text_w, h, text,
             font=FONT_MONO if mono else FONT_BODY, size=size,
             color=text_color, bold=True,
             anchor=MSO_ANCHOR.MIDDLE,
             align=(PP_ALIGN.CENTER if center else PP_ALIGN.LEFT),
             tracking=80)


def section_heading(slide, kicker_text, headline, lede=None, *,
                    kicker_color=GOLD, two_line=True, divider_color=GOLD):
    kicker(slide, MARGIN_L, KICKER_Y, CONTENT_W, kicker_text, color=kicker_color)
    h2_h = H2_H_2L if two_line else H2_H_1L
    add_text(slide, MARGIN_L, H2_Y, CONTENT_W, h2_h,
             headline, font=FONT_DISPLAY, size=(38 if two_line else 44),
             color=CREAM, italic=True, line_spacing=1.05)
    if lede:
        add_text(slide, MARGIN_L, LEDE_Y, Inches(9.6), LEDE_H,
                 lede, font=FONT_BODY, size=12.5, color=SOFT,
                 line_spacing=1.35)
    # accent divider — short gold + softer ghost extending
    add_rect(slide, MARGIN_L, DIV_Y, Inches(0.55), Inches(0.04),
             fill=divider_color, radius=0.5)
    ghost = add_rect(slide, MARGIN_L + Inches(0.65), DIV_Y + Inches(0.014),
                     Inches(2.4), Inches(0.012),
                     fill=RGBColor(0xFF, 0xFF, 0xFF), radius=0.5)
    shape_alpha(ghost, 18)


def footer(slide, page_num, total=12, sub="Investor deck · 2026"):
    brand_mark(slide, MARGIN_L, Inches(7.05), scale=0.62, sub=sub)
    # divider
    ln = slide.shapes.add_connector(1, MARGIN_L, FOOTER_DIV,
                                    SW - MARGIN_R, FOOTER_DIV)
    ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ln.line.width = Pt(0.5)
    line_alpha(ln, 8)
    # page indicator
    add_text(slide, SW - MARGIN_R - Inches(1.4), FOOTER_Y + Inches(0.03),
             Inches(1.4), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             font=FONT_MONO, size=10, color=MUTED,
             align=PP_ALIGN.RIGHT, bold=True, tracking=120)


# ===========================================================================
# SLIDE 01 — TITLE
# ===========================================================================
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('coral', 'teal', 'gold'))

    brand_mark(s, MARGIN_L, Inches(0.55), scale=1.0)

    # top-right tag
    chip(s, SW - MARGIN_R - Inches(2.5), Inches(0.62),
         Inches(2.5), Inches(0.42),
         "SEED · INVESTOR DECK", fill_color=PANEL,
         text_color=GOLD, border=STROKE_HI, center=True)

    # giant wordmark
    add_text(s, MARGIN_L, Inches(2.20), Inches(12.0), Inches(2.2),
             "Elyra", font=FONT_DISPLAY, size=148,
             color=CREAM, italic=True, line_spacing=1.0)

    # gradient sliver (cream→gold→coral via three blocks)
    bx = MARGIN_L + Inches(0.04); by = Inches(4.32)
    add_rect(s, bx,                 by, Inches(1.4), Inches(0.06),
             fill=GOLD, radius=0.5)
    bar2 = add_rect(s, bx + Inches(1.4),  by, Inches(1.6), Inches(0.06),
                    fill=CORAL, radius=0.5); shape_alpha(bar2, 92)
    bar3 = add_rect(s, bx + Inches(3.0),  by, Inches(0.9), Inches(0.06),
                    fill=VIOLET, radius=0.5); shape_alpha(bar3, 70)

    # tagline (mixed run)
    add_text(s, MARGIN_L, Inches(4.55), Inches(12.0), Inches(0.78),
             [
                 ("When identity is sensitive, ",
                  {'italic': True, 'color': SOFT,
                   'font': FONT_DISPLAY, 'size': 30}),
                 ("trust",
                  {'italic': True, 'color': GOLD,
                   'font': FONT_DISPLAY, 'size': 30}),
                 (" becomes the product.",
                  {'italic': True, 'color': SOFT,
                   'font': FONT_DISPLAY, 'size': 30}),
             ], line_spacing=1.05)

    # subtitle
    add_text(s, MARGIN_L, Inches(5.32), Inches(11.0), Inches(0.36),
             "PEHCHAAN LAYER  ·  INDIA-FIRST  ·  PRIVACY-FIRST LGBTQIA+ CONNECTION",
             font=FONT_BODY, size=10.5, color=MUTED, bold=True, tracking=240)

    # proof row — 4 cards
    proof = [
        ("256-bit", "AES private layer",  GOLD),
        ("384d",    "intent embeddings",  TEAL),
        ("₹499+",   "premium ARPU path",  CORAL),
        ("91",      "trust score, live",  VIOLET),
    ]
    px = MARGIN_L; py = Inches(6.05)
    pw = Inches(2.91); ph = Inches(0.80); gap = Inches(0.13)
    for i, (val, label, c) in enumerate(proof):
        cx = px + (pw + gap) * i
        glass_card(s, cx, py, pw, ph, radius=0.22)
        add_text(s, cx + Inches(0.22), py + Inches(0.10),
                 pw - Inches(0.55), Inches(0.45),
                 val, font=FONT_DISPLAY, size=22,
                 color=CREAM, italic=True)
        add_text(s, cx + Inches(0.22), py + Inches(0.50),
                 pw - Inches(0.55), Inches(0.30),
                 label, font=FONT_BODY, size=9.5,
                 color=MUTED, bold=True, tracking=80)
        add_ellipse(s, cx + pw - Inches(0.36), py + Inches(0.24),
                    Inches(0.14), Inches(0.14), c)

    # discreet page indicator
    add_text(s, SW - MARGIN_R - Inches(1.4), FOOTER_Y + Inches(0.03),
             Inches(1.4), Inches(0.3),
             "01 / 12", font=FONT_MONO, size=10, color=MUTED,
             align=PP_ALIGN.RIGHT, bold=True, tracking=120)


# ===========================================================================
# SLIDE 02 — PROBLEM
# ===========================================================================
def slide_problem():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('coral', 'gold'))
    section_heading(
        s, "01 · The Problem",
        "Dating products assume\neveryone can be public.",
        "For most LGBTQIA+ users in India, disclosure, harassment, fake profiles, "
        "and unsafe offline meetings aren't edge cases — they're core adoption "
        "barriers no swipe-first product solves.",
        kicker_color=CORAL, divider_color=CORAL,
    )

    cards = [
        ("73%",    "fear public exposure",
         "Identity-sensitive users won't post photos, names, or location on swipe apps that treat profiles as public-by-default.",
         CORAL),
        ("1 in 4", "encounter fake profiles",
         "Without verification, moderation, and risk scoring, trust collapses before the first conversation begins.",
         GOLD),
        ("0",      "safety nets for offline",
         "First meetings happen offline — but no major app builds emergency contacts, check-ins, or SOS into the product loop.",
         VIOLET),
    ]
    cx = MARGIN_L; cy = BODY_Y
    cw = Inches(3.85); ch = Inches(3.10); gap = Inches(0.14)
    for i, (stat, head, body, c) in enumerate(cards):
        x = cx + (cw + gap) * i
        glass_card(s, x, cy, cw, ch, radius=0.07)
        # left accent
        add_rect(s, x + Inches(0.34), cy + Inches(0.34),
                 Inches(0.045), Inches(2.42), fill=c, radius=0.5)
        add_text(s, x + Inches(0.55), cy + Inches(0.30),
                 cw - Inches(0.85), Inches(1.30),
                 stat, font=FONT_DISPLAY, size=66,
                 color=CREAM, italic=True, line_spacing=1.0)
        add_text(s, x + Inches(0.55), cy + Inches(1.65),
                 cw - Inches(0.85), Inches(0.34),
                 head, font=FONT_BODY, size=13.5,
                 color=c, bold=True, tracking=20)
        add_text(s, x + Inches(0.55), cy + Inches(2.05),
                 cw - Inches(0.85), Inches(0.95),
                 body, font=FONT_BODY, size=10.5, color=SOFT,
                 line_spacing=1.35)

    footer(s, 2)


# ===========================================================================
# SLIDE 03 — WHY NOW
# ===========================================================================
def slide_why_now():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('teal', 'violet'))
    section_heading(
        s, "02 · Why Now",
        "Three tailwinds converge in 2026.",
        "Cultural shift, real willingness-to-pay for privacy, and AI safety "
        "infrastructure that's finally cheap enough to ship as a default.",
        kicker_color=TEAL, two_line=False, divider_color=TEAL,
    )

    pillars = [
        ("01", "Identity is mainstream",
         "Pride, queer creators, and Section 377's reversal moved LGBTQIA+ identity from invisible to digitally native — but product hasn't caught up.",
         CORAL),
        ("02", "Privacy is a paid behavior",
         "Indian users pay for VPNs, vault apps, and incognito tiers. Willingness-to-pay for privacy as a feature is at an all-time high.",
         GOLD),
        ("03", "AI safety is finally cheap",
         "Open-weights moderation, embeddings, and on-device blur are production-grade. What cost $50M five years ago is now a fraction.",
         TEAL),
    ]
    cx = MARGIN_L; cy = BODY_Y
    cw = Inches(3.85); ch = Inches(3.10); gap = Inches(0.14)
    for i, (num, head, body, c) in enumerate(pillars):
        x = cx + (cw + gap) * i
        glass_card(s, x, cy, cw, ch, radius=0.08)
        add_text(s, x + Inches(0.36), cy + Inches(0.32),
                 Inches(1.0), Inches(0.5),
                 num, font=FONT_DISPLAY, size=28,
                 color=c, italic=True)
        # divider
        ln = s.shapes.add_connector(1, x + Inches(0.36), cy + Inches(0.92),
                                    x + cw - Inches(0.36), cy + Inches(0.92))
        ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ln.line.width = Pt(0.5); line_alpha(ln, 14)
        add_text(s, x + Inches(0.36), cy + Inches(1.10),
                 cw - Inches(0.6), Inches(0.7),
                 head, font=FONT_DISPLAY, size=22,
                 color=CREAM, italic=True, line_spacing=1.05)
        add_text(s, x + Inches(0.36), cy + Inches(1.85),
                 cw - Inches(0.6), Inches(1.15),
                 body, font=FONT_BODY, size=11, color=SOFT,
                 line_spacing=1.35)

    footer(s, 3, sub="Investor deck · 2026")


# ===========================================================================
# SLIDE 04 — SOLUTION
# ===========================================================================
def slide_solution():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('teal', 'coral', 'gold'))
    section_heading(
        s, "03 · The Solution",
        "Elyra is the Trust OS\nfor modern connection.",
        "Pehchaan Layer dual identity, explainable matching, AI moderation, "
        "risk scoring, Safe Date, and premium privacy — one product surface.",
        kicker_color=GOLD,
    )

    feats = [
        ("01", "Dual identity",
         "Public profile stays light. Private fields stay AES-256 encrypted, reveal-based, and audit-logged.", CORAL),
        ("02", "Intent matching",
         "No swipe pressure. Candidates ranked by intent, preference, distance, and 384d embeddings.", GOLD),
        ("03", "AI safety",
         "Toxicity scoring, fake-profile heuristics, image checks, and explainable risk traces in-product.", TEAL),
        ("04", "Safe Date",
         "Emergency contacts, timed check-ins, live location, and SOS — designed into the product loop.", VIOLET),
    ]
    cx = MARGIN_L; cy = BODY_Y
    cw = Inches(2.84); ch = Inches(3.10); gap = Inches(0.13)
    for i, (num, head, body, c) in enumerate(feats):
        x = cx + (cw + gap) * i
        glass_card(s, x, cy, cw, ch, radius=0.08)
        add_text(s, x + Inches(0.32), cy + Inches(0.32),
                 cw - Inches(0.6), Inches(0.6),
                 num, font=FONT_DISPLAY, size=30,
                 color=c, italic=True)
        # divider
        ln = s.shapes.add_connector(1, x + Inches(0.32), cy + Inches(1.00),
                                    x + cw - Inches(0.32), cy + Inches(1.00))
        ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); ln.line.width = Pt(0.5)
        line_alpha(ln, 14)
        add_text(s, x + Inches(0.32), cy + Inches(1.16),
                 cw - Inches(0.6), Inches(0.5),
                 head, font=FONT_DISPLAY, size=21,
                 color=CREAM, italic=True)
        add_text(s, x + Inches(0.32), cy + Inches(1.78),
                 cw - Inches(0.6), Inches(1.20),
                 body, font=FONT_BODY, size=10.5, color=SOFT,
                 line_spacing=1.35)
        add_ellipse(s, x + cw - Inches(0.52), cy + Inches(0.40),
                    Inches(0.18), Inches(0.18), c)

    footer(s, 4)


# ===========================================================================
# SLIDE 05 — DEMO MAP
# ===========================================================================
def slide_demo_map():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('teal', 'gold'))
    section_heading(
        s, "04 · Interactive Product Demo",
        "The full investor flow,\nin under 90 seconds.",
        "Five panels, one cockpit: simulate onboarding, dual-identity reveal, "
        "match generation, AI moderation, and Safe Date activation.",
        kicker_color=TEAL, divider_color=TEAL,
    )

    # left rail
    rx = MARGIN_L; ry = BODY_Y
    rw = Inches(4.4); rh = Inches(3.10)
    glass_card(s, rx, ry, rw, rh, radius=0.06)
    add_text(s, rx + Inches(0.36), ry + Inches(0.30),
             Inches(2.6), Inches(0.3),
             "INVESTOR FLOW", font=FONT_BODY, size=9,
             color=MUTED, bold=True, tracking=240)
    add_text(s, rx + rw - Inches(1.2), ry + Inches(0.30),
             Inches(0.84), Inches(0.3),
             "1 / 5", font=FONT_MONO, size=10, color=GOLD,
             bold=True, align=PP_ALIGN.RIGHT)

    rail = [
        ("1", "Onboarding",     "Identity + intent capture", CORAL),
        ("2", "Dual identity",  "Pehchaan Layer reveal",     GOLD),
        ("3", "Matching AI",    "Explainable scoring",       TEAL),
        ("4", "Moderated chat", "Visible AI decisions",      VIOLET),
        ("5", "Safe Date",      "Offline safety loop",       BLUE),
    ]
    sy = ry + Inches(0.78)
    for n, head, sub, c in rail:
        sx = rx + Inches(0.36)
        d = Inches(0.34)
        add_ellipse(s, sx, sy + Inches(0.04), d, d, c)
        add_text(s, sx, sy + Inches(0.04), d, d, n,
                 font=FONT_BODY, size=11, color=INK_TEXT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sx + Inches(0.52), sy,
                 rw - Inches(1.0), Inches(0.26),
                 head, font=FONT_BODY, size=12.5, color=CREAM, bold=True)
        add_text(s, sx + Inches(0.52), sy + Inches(0.24),
                 rw - Inches(1.0), Inches(0.24),
                 sub, font=FONT_BODY, size=9.5, color=MUTED, tracking=40)
        sy += Inches(0.45)

    # right panel — explainability mock
    px = MARGIN_L + rw + Inches(0.18); py = BODY_Y
    pw = SW - MARGIN_R - px; ph = Inches(3.10)
    glass_card(s, px, py, pw, ph, radius=0.06)
    add_text(s, px + Inches(0.42), py + Inches(0.30),
             Inches(4.0), Inches(0.3),
             "AI EXPLAINABILITY · LIVE", font=FONT_BODY, size=9,
             color=GOLD, bold=True, tracking=240)
    # status pill
    chip(s, px + pw - Inches(1.45), py + Inches(0.27),
         Inches(1.05), Inches(0.34),
         "TRUST 91", fill_color=PANEL_2, text_color=TEAL,
         border=STROKE, dot_color=TEAL, size=9, center=False)
    add_text(s, px + Inches(0.42), py + Inches(0.66),
             pw - Inches(0.84), Inches(0.55),
             "Match · Aarav, 28  ·  92%",
             font=FONT_DISPLAY, size=24, color=CREAM, italic=True)

    scores = [
        ("Intent match",         0.96, "30%", TEAL),
        ("Embedding similarity", 0.88, "35%", BLUE),
        ("Distance score",       0.74, "20%", GOLD),
        ("Preference fit",       0.91, "15%", CORAL),
    ]
    sy = py + Inches(1.40)
    for label, frac, weight, c in scores:
        add_text(s, px + Inches(0.42), sy,
                 Inches(2.4), Inches(0.28),
                 label, font=FONT_BODY, size=10.5, color=SOFT)
        add_text(s, px + Inches(2.85), sy,
                 Inches(0.6), Inches(0.28),
                 weight, font=FONT_MONO, size=10, color=GOLD, bold=True)
        track_x = px + Inches(3.55); track_w = pw - Inches(4.0)
        add_rect(s, track_x, sy + Inches(0.10),
                 track_w, Inches(0.10), fill=PANEL_3,
                 radius=0.5, alpha=70)
        add_rect(s, track_x, sy + Inches(0.10),
                 int(track_w * frac), Inches(0.10), fill=c, radius=0.5)
        sy += Inches(0.38)

    footer(s, 5)


# ===========================================================================
# SLIDE 06 — CORE FEATURES
# ===========================================================================
def slide_core_features():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('coral', 'violet'))
    section_heading(
        s, "05 · Core Features",
        "The moat, in plain English.",
        "Each is a feature competitors could ship. The moat is shipping them "
        "together, visibly, and as the product surface itself.",
        kicker_color=VIOLET, two_line=False, divider_color=VIOLET,
    )

    rows = [
        ("Pehchaan Layer", "Dual identity — AES-256-GCM, reveal-by-consent, audit-trailed, revocable. Architecture, not a setting.", CORAL),
        ("Explainable matching", "Composite score with weights surfaced to the user. Builds trust and satisfies algorithmic-transparency regulation.", GOLD),
        ("Visible AI moderation", "Every chat decision shows category, toxicity, action, and a reversible \"View anyway\" path.", TEAL),
        ("Safe Date as product", "Emergency contact, check-ins, SOS, live route, missed-check-in escalation — first-class flow.", VIOLET),
        ("Composite trust score", "Verification + toxicity + reports + fake-profile + age. One number that earns the right to charge.", BLUE),
    ]
    rx = MARGIN_L; ry = BODY_Y
    rw = CONTENT_W; rh = Inches(0.56); gap = Inches(0.10)
    for i, (head, body, c) in enumerate(rows):
        y = ry + (rh + gap) * i
        glass_card(s, rx, y, rw, rh, radius=0.30)
        # accent capsule on left
        add_rect(s, rx + Inches(0.30), y + Inches(0.18),
                 Inches(0.20), Inches(0.20), fill=c, radius=0.5)
        add_text(s, rx + Inches(0.78), y, Inches(3.6), rh,
                 head, font=FONT_BODY, size=12.5,
                 color=CREAM, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        # vertical divider
        sep = s.shapes.add_connector(1, rx + Inches(4.45), y + Inches(0.13),
                                     rx + Inches(4.45), y + rh - Inches(0.13))
        sep.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); sep.line.width = Pt(0.5)
        line_alpha(sep, 14)
        add_text(s, rx + Inches(4.65), y, rw - Inches(4.85), rh,
                 body, font=FONT_BODY, size=11,
                 color=SOFT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 6)


# ===========================================================================
# SLIDE 07 — DIFFERENTIATION
# ===========================================================================
def slide_diff():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('coral', 'teal'))
    section_heading(
        s, "06 · Differentiation",
        "Why we win — not copyable\nin a quarter.",
        "Trust is a system, not a feature. Replicating dual identity + AI safety + "
        "Safe Date + audit infrastructure is a 12–18 month rebuild — and it "
        "cannibalizes incumbents' existing products.",
        kicker_color=CORAL, divider_color=CORAL,
    )

    headers = ["Vector", "Swipe apps", "Niche queer apps", "Elyra"]
    rows = [
        ("Public-by-default",  "Yes (assumed)",   "Yes (mostly)",   "No — Pehchaan Layer"),
        ("AI moderation",      "Off-screen",      "Light / none",   "Visible & reversible"),
        ("Offline safety",     "Buried in help",  "Absent",         "Safe Date — core flow"),
        ("Identity model",     "Single profile",  "Single profile", "Dual, encrypted, audited"),
        ("Matching",           "Swipe queue",     "Swipe queue",    "No-swipe, intent-first"),
        ("India queer fit",    "Generic",         "Imported UX",    "India-first, vernacular"),
    ]
    tx = MARGIN_L; ty = BODY_Y
    cols = [Inches(2.55), Inches(2.55), Inches(2.85), Inches(3.88)]
    rh = Inches(0.40); gap = Inches(0.05)

    # header row
    cx = tx
    for i, (h, w) in enumerate(zip(headers, cols)):
        is_elyra = (i == 3)
        bg_color = GOLD if is_elyra else PANEL_2
        text_color = INK_TEXT if is_elyra else GOLD
        add_rect(s, cx, ty, w, Inches(0.46),
                 fill=bg_color, line=STROKE, line_w=0.5,
                 radius=0.30, alpha=(100 if is_elyra else 88),
                 line_alpha_pct=70)
        add_text(s, cx + Inches(0.30), ty,
                 w - Inches(0.6), Inches(0.46),
                 h.upper(), font=FONT_BODY, size=9.5,
                 color=text_color, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, tracking=160)
        cx += w
    ty += Inches(0.46) + gap

    for r_idx, row in enumerate(rows):
        cx = tx
        for i, (val, w) in enumerate(zip(row, cols)):
            is_elyra = (i == 3)
            fill = PANEL_3 if is_elyra else PANEL
            alpha = 92 if is_elyra else 80
            add_rect(s, cx, ty, w, rh,
                     fill=fill, line=STROKE_SOFT, line_w=0.4,
                     radius=0.25, alpha=alpha, line_alpha_pct=80)
            if is_elyra:
                # small left dot
                add_ellipse(s, cx + Inches(0.18), ty + rh / 2 - Inches(0.07),
                            Inches(0.14), Inches(0.14), GOLD)
                text_x = cx + Inches(0.42)
            else:
                text_x = cx + Inches(0.30)
            add_text(s, text_x, ty,
                     w - (text_x - cx) - Inches(0.20), rh,
                     val, font=FONT_BODY, size=10.5,
                     color=(GOLD if is_elyra else SOFT),
                     bold=is_elyra, anchor=MSO_ANCHOR.MIDDLE)
            cx += w
        ty += rh + gap

    footer(s, 7)


# ===========================================================================
# SLIDE 08 — MARKET
# ===========================================================================
def slide_market():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('teal', 'gold', 'violet'))
    section_heading(
        s, "07 · Market Opportunity",
        "India today.\nThe world tomorrow.",
        "Start where identity-sensitive connection has acute need; expand into "
        "every market where trust is a differentiating product moat.",
        kicker_color=GOLD,
    )

    stats = [
        ("$2.5B", "India dating TAM by 2028",       GOLD),
        ("2.5M+", "LGBTQIA+ digital natives",        CORAL),
        ("6×",    "privacy-paid willingness-to-pay", TEAL),
    ]
    sx = MARGIN_L; sy = BODY_Y
    sw = Inches(2.10); sh = Inches(3.10); gap = Inches(0.14)
    for i, (val, label, c) in enumerate(stats):
        x = sx + (sw + gap) * i
        glass_card(s, x, sy, sw, sh, radius=0.08)
        add_text(s, x + Inches(0.28), sy + Inches(0.55),
                 sw - Inches(0.5), Inches(1.4),
                 val, font=FONT_DISPLAY, size=44,
                 color=CREAM, italic=True, line_spacing=1.0)
        add_rect(s, x + Inches(0.28), sy + Inches(1.95),
                 Inches(1.0), Inches(0.05), fill=c, radius=0.5)
        add_text(s, x + Inches(0.28), sy + Inches(2.15),
                 sw - Inches(0.5), Inches(0.85),
                 label, font=FONT_BODY, size=10.5,
                 color=SOFT, line_spacing=1.30)

    # right: phased thesis
    px = sx + (sw + gap) * 3 + Inches(0.10); py = BODY_Y
    pw = SW - MARGIN_R - px; ph = Inches(3.10)
    glass_card(s, px, py, pw, ph, radius=0.06)
    kicker(s, px + Inches(0.36), py + Inches(0.30),
           Inches(4), "EXPANSION THESIS", color=TEAL, size=9.5)

    phases = [
        ("Phase 1 · 2026",  "India web + mobile launch",                  TEAL),
        ("Phase 2 · 2027",  "SEA + MENA · privacy as survival feature",   GOLD),
        ("Phase 3 · 2028+", "EU + LATAM · privacy as regulated standard", CORAL),
    ]
    py2 = py + Inches(0.86)
    for head, body, c in phases:
        add_ellipse(s, px + Inches(0.36), py2 + Inches(0.18),
                    Inches(0.18), Inches(0.18), c)
        add_text(s, px + Inches(0.72), py2,
                 pw - Inches(1.0), Inches(0.30),
                 head, font=FONT_BODY, size=12, color=CREAM, bold=True)
        add_text(s, px + Inches(0.72), py2 + Inches(0.30),
                 pw - Inches(1.0), Inches(0.34),
                 body, font=FONT_BODY, size=10.5, color=SOFT)
        py2 += Inches(0.74)

    footer(s, 8)


# ===========================================================================
# SLIDE 09 — BUSINESS MODEL
# ===========================================================================
def slide_business():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('gold', 'coral'))
    section_heading(
        s, "08 · Business Model",
        "Privacy creates a premium\nwillingness-to-pay loop.",
        "Free for adoption. Paid for power. We charge for boost, incognito, "
        "filters, and AI starters — never for safety.",
        kicker_color=GOLD,
    )

    plans = [
        ("Free",    "₹0",    ["Basic intent match", "5 likes / day", "Trust score visible"],     False),
        ("Plus",    "₹499",  ["Unlimited likes", "See who liked you", "No ads · priority mod"],   True),
        ("Premium", "₹999",  ["Advanced filters", "Priority matching", "LLM conversation starters"], False),
        ("Elite",   "₹1999", ["Profile boost", "Incognito mode", "Priority human support"],      False),
    ]
    cx = MARGIN_L; cy = BODY_Y
    cw = Inches(2.84); ch = Inches(3.05); gap = Inches(0.13)

    for i, (name, price, lines, featured) in enumerate(plans):
        x = cx + (cw + gap) * i
        if featured:
            # outer glow halo
            halo = add_rect(s, x - Inches(0.10), cy - Inches(0.10),
                            cw + Inches(0.20), ch + Inches(0.20),
                            fill=GOLD, radius=0.10, alpha=18)
            # featured solid card
            add_rect(s, x, cy, cw, ch, fill=GOLD, line=GOLD,
                     line_w=1.0, radius=0.10, alpha=100,
                     line_alpha_pct=100)
            tag_color = INK_TEXT
            txt_color = INK_TEXT
            sub_color = INK_TEXT
        else:
            glass_card(s, x, cy, cw, ch, radius=0.10)
            tag_color = GOLD
            txt_color = CREAM
            sub_color = SOFT

        # tag
        tag_text = "MOST POPULAR" if featured else name.upper()
        add_text(s, x + Inches(0.36), cy + Inches(0.32),
                 cw - Inches(0.6), Inches(0.30),
                 tag_text, font=FONT_BODY, size=9, color=tag_color,
                 bold=True, tracking=240)
        # plan name (only if not featured — featured uses tag for it differently)
        add_text(s, x + Inches(0.36), cy + Inches(0.66),
                 cw - Inches(0.6), Inches(0.36),
                 name, font=FONT_BODY, size=12,
                 color=tag_color, bold=True)
        # price
        add_text(s, x + Inches(0.36), cy + Inches(1.04),
                 cw - Inches(0.6), Inches(0.95),
                 price, font=FONT_DISPLAY, size=46,
                 color=txt_color, italic=True, line_spacing=1.0)

        # divider
        ln = s.shapes.add_connector(1, x + Inches(0.36), cy + Inches(2.00),
                                    x + cw - Inches(0.36), cy + Inches(2.00))
        if featured:
            ln.line.color.rgb = INK_TEXT
            ln.line.width = Pt(0.5)
            line_alpha(ln, 35)
        else:
            ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            ln.line.width = Pt(0.5)
            line_alpha(ln, 14)

        # bullet lines
        ly = cy + Inches(2.13)
        for ln_text in lines:
            add_ellipse(s, x + Inches(0.36), ly + Inches(0.10),
                        Inches(0.07), Inches(0.07),
                        (INK_TEXT if featured else GOLD))
            add_text(s, x + Inches(0.55), ly,
                     cw - Inches(0.85), Inches(0.26),
                     ln_text, font=FONT_BODY, size=10,
                     color=sub_color)
            ly += Inches(0.27)

    # bottom unit-econ strip (sits above the footer divider)
    add_text(s, MARGIN_L, Inches(6.78), CONTENT_W, Inches(0.24),
             "TARGET ARPU ₹600 / MO   ·   6–8% PAYING CONVERSION   ·   "
             "CAC PAYBACK ≤ 4 MONTHS   ·   ANNUAL SAVE 17%",
             font=FONT_BODY, size=9, color=MUTED,
             bold=True, tracking=200, align=PP_ALIGN.LEFT)

    footer(s, 9)


# ===========================================================================
# SLIDE 10 — GO-TO-MARKET
# ===========================================================================
def slide_gtm():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('teal', 'coral'))
    section_heading(
        s, "09 · Go-to-Market",
        "Land community-first.\nConvert through trust.",
        "Earned reach > paid reach in identity-sensitive communities. "
        "Our trust mechanics double as a self-credentialing growth loop.",
        kicker_color=TEAL, divider_color=TEAL,
    )

    steps = [
        ("M0–M3",  "Partnerships",       "Queer creators, pride orgs, campus alliances · 8 metro & tier-1 cities", CORAL),
        ("M3–M6",  "Verified seed loop", "Verified invites raise composite trust score · gamified, privacy-first growth", GOLD),
        ("M6–M9",  "Vernacular content", "Hindi · Tamil · Bengali content on consent, safe meeting, identity",      TEAL),
        ("M9–M12", "Mobile + paid",      "Expo mobile launch · privacy-led creative · Premium / Elite upsells",     VIOLET),
        ("M12+",   "B2B venues",         "Safe-Date-friendly cafes & bars · network effect · earned discovery",     BLUE),
    ]
    sx = MARGIN_L; sy = BODY_Y
    n = len(steps)
    gap = Inches(0.10)
    sw = (CONTENT_W - gap * (n - 1)) / n
    sh = Inches(3.10)

    for i, (window, head, body, c) in enumerate(steps):
        x = sx + (sw + gap) * i
        glass_card(s, x, sy, sw, sh, radius=0.08)
        # window pill
        pill_w = sw - Inches(0.72)
        add_rect(s, x + Inches(0.36), sy + Inches(0.32),
                 pill_w, Inches(0.32),
                 fill=c, radius=0.5, alpha=20)
        add_text(s, x + Inches(0.36), sy + Inches(0.32),
                 pill_w, Inches(0.32),
                 window, font=FONT_MONO, size=9.5, color=c,
                 bold=True, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER, tracking=80)
        add_text(s, x + Inches(0.36), sy + Inches(0.85),
                 sw - Inches(0.72), Inches(0.6),
                 head, font=FONT_DISPLAY, size=18,
                 color=CREAM, italic=True, line_spacing=1.05)
        # divider
        ln = s.shapes.add_connector(1, x + Inches(0.36), sy + Inches(1.50),
                                    x + sw - Inches(0.36), sy + Inches(1.50))
        ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); ln.line.width = Pt(0.5)
        line_alpha(ln, 14)
        add_text(s, x + Inches(0.36), sy + Inches(1.62),
                 sw - Inches(0.72), Inches(1.40),
                 body, font=FONT_BODY, size=10, color=SOFT,
                 line_spacing=1.35)
        # connector arrow between cards
        if i < n - 1:
            arrow_y = sy + sh / 2
            arrow = s.shapes.add_connector(1, x + sw + Inches(0.005),
                                           arrow_y,
                                           x + sw + gap - Inches(0.005),
                                           arrow_y)
            arrow.line.color.rgb = STROKE_HI
            arrow.line.width = Pt(0.75)

    footer(s, 10)


# ===========================================================================
# SLIDE 11 — VISION
# ===========================================================================
def slide_vision():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('violet', 'gold'))
    section_heading(
        s, "10 · Vision",
        "From queer dating in India\nto the world's Trust OS.",
        "The dating app is the wedge. The prize is being the trust infrastructure "
        "for every product where identity is sensitive.",
        kicker_color=VIOLET,
    )

    horizons = [
        ("Y1",   "India launch",
         "Web + mobile. Pehchaan Layer becomes a household phrase in queer urban India.", CORAL),
        ("Y2–3", "SEA & MENA",
         "Dual identity becomes the standard for identity-sensitive markets globally.",   GOLD),
        ("Y3–5", "Trust Layer platform",
         "Pehchaan APIs power identity-sensitive flows in fintech, healthcare, hiring.",  TEAL),
        ("Y5+",  "Global default",
         "Elyra is to identity-safe connection what Stripe became to payments.",          VIOLET),
    ]

    # equalised cards · subtle stair conveyed via a horizontal accent bar that
    # grows in width with each horizon (left→right ramps the bar).
    n = len(horizons)
    gap = Inches(0.13)
    cw = (CONTENT_W - gap * (n - 1)) / n
    ch = Inches(3.10)
    base_y = BODY_Y

    for i, (era, head, body, c) in enumerate(horizons):
        x = MARGIN_L + (cw + gap) * i
        y = base_y
        glass_card(s, x, y, cw, ch, radius=0.08)
        # era chip
        add_rect(s, x + Inches(0.36), y + Inches(0.34),
                 Inches(1.0), Inches(0.34),
                 fill=c, radius=0.5, alpha=20)
        add_text(s, x + Inches(0.36), y + Inches(0.34),
                 Inches(1.0), Inches(0.34), era,
                 font=FONT_MONO, size=10, color=c, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.36), y + Inches(0.86),
                 cw - Inches(0.72), Inches(0.66),
                 head, font=FONT_DISPLAY, size=20,
                 color=CREAM, italic=True, line_spacing=1.05)
        add_text(s, x + Inches(0.36), y + Inches(1.62),
                 cw - Inches(0.72), Inches(1.05),
                 body, font=FONT_BODY, size=10.5, color=SOFT,
                 line_spacing=1.35)
        # bottom accent bar grows by horizon index — visual stair
        bar_w = Inches(0.6) + Inches(0.55) * i
        add_rect(s, x + Inches(0.36), y + ch - Inches(0.30),
                 bar_w, Inches(0.045), fill=c, radius=0.5)

    footer(s, 11)


# ===========================================================================
# SLIDE 12 — THE ASK
# ===========================================================================
def slide_ask():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, orbs=('coral', 'teal', 'gold'))
    section_heading(
        s, "11 · The Ask",
        "Trust is the next dating moat.\nElyra is building it.",
        "We're raising seed to launch India, build the AI safety stack, and ship "
        "mobile. Investors get a privacy-defensible platform aimed at an "
        "underserved, high-LTV community.",
        kicker_color=GOLD,
    )

    # left: Use of funds
    lx = MARGIN_L; ly = BODY_Y
    lw = Inches(6.30); lh = Inches(3.10)
    glass_card(s, lx, ly, lw, lh, radius=0.06)
    kicker(s, lx + Inches(0.36), ly + Inches(0.32), Inches(4),
           "USE OF FUNDS · 18 MONTHS", color=GOLD, size=9.5)

    fund = [
        (40, "Engineering",       "Web v1, mobile v1, FastAPI · PG · Mongo · pgvector on K8s",  CORAL),
        (25, "AI safety stack",   "Toxicity, fake-profile, image mod, reasoning UX, review queue", TEAL),
        (20, "GTM",               "Creator partnerships, vernacular content, paid mobile launch",  GOLD),
        (10, "Trust & compliance","Legal · DPDP-ready · SOC 2 path · on-call moderation",          VIOLET),
        ( 5, "Reserve",           "Buffer for hiring & opportunistic spend",                        BLUE),
    ]
    fy = ly + Inches(0.80)
    for pct, head, body, c in fund:
        # percentage
        add_text(s, lx + Inches(0.36), fy - Inches(0.02),
                 Inches(0.95), Inches(0.40),
                 f"{pct}%", font=FONT_DISPLAY, size=20,
                 color=c, italic=True)
        # heading row + bar
        bar_x = lx + Inches(1.40); bar_w = lw - Inches(1.78)
        add_text(s, bar_x, fy - Inches(0.04),
                 bar_w, Inches(0.22),
                 head, font=FONT_BODY, size=10.5,
                 color=CREAM, bold=True)
        # bar bg
        add_rect(s, bar_x, fy + Inches(0.18),
                 bar_w, Inches(0.10),
                 fill=PANEL_3, radius=0.5, alpha=70)
        # bar fill (40% scales to full bar)
        add_rect(s, bar_x, fy + Inches(0.18),
                 int(bar_w * (pct / 40)), Inches(0.10),
                 fill=c, radius=0.5)
        # subtitle
        add_text(s, bar_x, fy + Inches(0.30),
                 bar_w, Inches(0.20),
                 body, font=FONT_BODY, size=9, color=MUTED)
        fy += Inches(0.42)

    # right: milestones
    rx = lx + lw + Inches(0.18); ry = BODY_Y
    rw = SW - MARGIN_R - rx; rh = Inches(3.10)
    glass_card(s, rx, ry, rw, rh, radius=0.06)
    kicker(s, rx + Inches(0.36), ry + Inches(0.32), Inches(4),
           "MILESTONES · 18 MONTHS", color=CORAL, size=9.5)

    miles = [
        ("Q1 '26", "India web launch",   "25k seed users",         TEAL),
        ("Q3 '26", "Mobile launch",      "Paying tier live",       GOLD),
        ("Q1 '27", "250k MAU",           "6%+ paying · Series A",  CORAL),
        ("Q3 '27", "SEA / MENA pilot",   "First international",    VIOLET),
    ]
    my = ry + Inches(0.85)
    for q, head, sub, c in miles:
        # quarter chip
        add_rect(s, rx + Inches(0.36), my,
                 Inches(1.0), Inches(0.36),
                 fill=c, radius=0.5, alpha=20)
        add_text(s, rx + Inches(0.36), my,
                 Inches(1.0), Inches(0.36),
                 q, font=FONT_MONO, size=10,
                 color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(1.55), my - Inches(0.02),
                 rw - Inches(1.85), Inches(0.24),
                 head, font=FONT_BODY, size=12,
                 color=CREAM, bold=True)
        add_text(s, rx + Inches(1.55), my + Inches(0.20),
                 rw - Inches(1.85), Inches(0.22),
                 sub, font=FONT_BODY, size=10, color=MUTED)
        my += Inches(0.50)

    # closing line (sits above footer divider)
    add_text(s, MARGIN_L, Inches(6.78), CONTENT_W, Inches(0.24),
             "WHEN IDENTITY IS SENSITIVE  —  TRUST BECOMES THE PRODUCT.",
             font=FONT_BODY, size=10.5, color=GOLD,
             bold=True, tracking=240)

    footer(s, 12)


# ---------------------------------------------------------------------------
# BUILD
# ---------------------------------------------------------------------------
slide_title()
slide_problem()
slide_why_now()
slide_solution()
slide_demo_map()
slide_core_features()
slide_diff()
slide_market()
slide_business()
slide_gtm()
slide_vision()
slide_ask()

OUT = "/home/kaarthikeya/Downloads/Elyra-main/Elyra-main/POC/PPT/Elyra_Investor_Pitch.pptx"
prs.save(OUT)
print("Saved:", OUT)
